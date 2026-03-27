import pytest
from pathlib import Path
from pptx import Presentation
from vbook.output.pptx_generator import generate_pptx


def test_generate_pptx_basic(tmp_path):
    """Test basic PPT generation with title and sections."""
    analysis = {
        "title": "测试演示",
        "outline": [
            {"title": "章节1", "summary": "这是第一章节的摘要"},
            {"title": "章节2", "summary": "这是第二章节的摘要"},
        ],
        "keywords": ["关键词1", "关键词2", "关键词3"],
    }

    assets_dir = tmp_path / "assets"
    assets_dir.mkdir()
    output_path = tmp_path / "test.pptx"

    generate_pptx(analysis, assets_dir, output_path)

    assert output_path.exists()
    prs = Presentation(str(output_path))
    # Title slide + 2 section slides + keywords slide = 4 slides
    assert len(prs.slides) == 4


def test_generate_pptx_with_screenshots(tmp_path):
    """Test PPT generation with screenshots."""
    # Create a dummy image
    from PIL import Image
    assets_dir = tmp_path / "assets"
    assets_dir.mkdir()
    img_path = assets_dir / "test_img.jpg"
    img = Image.new("RGB", (100, 100), color="red")
    img.save(str(img_path))

    analysis = {
        "title": "带图片的演示",
        "outline": [
            {
                "title": "章节1",
                "summary": "包含截图的章节",
                "screenshots": ["test_img.jpg"],
            }
        ],
        "keywords": ["测试"],
    }

    output_path = tmp_path / "test_with_img.pptx"
    generate_pptx(analysis, assets_dir, output_path)

    assert output_path.exists()
    prs = Presentation(str(output_path))
    # Title + 1 section + keywords = 3 slides
    assert len(prs.slides) == 3


def test_generate_pptx_no_keywords(tmp_path):
    """Test PPT generation without keywords."""
    analysis = {
        "title": "无关键词演示",
        "outline": [{"title": "章节1", "summary": "摘要"}],
        "keywords": [],
    }

    assets_dir = tmp_path / "assets"
    assets_dir.mkdir()
    output_path = tmp_path / "test_no_kw.pptx"

    generate_pptx(analysis, assets_dir, output_path)

    assert output_path.exists()
    prs = Presentation(str(output_path))
    # Title + 1 section + keywords (empty) = 3 slides
    assert len(prs.slides) == 3

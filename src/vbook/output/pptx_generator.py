# src/vbook/output/pptx_generator.py
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def generate_pptx(analysis: dict, assets_dir: Path, output_path: Path):
    """Generate PowerPoint presentation from analysis data.

    Args:
        analysis: Analysis data with title, outline, keywords
        assets_dir: Directory containing screenshot assets
        output_path: Path to save the .pptx file
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = analysis.get("title", "视频分析")
    keywords = analysis.get("keywords", [])
    subtitle.text = "关键词: " + ", ".join(keywords[:5]) if keywords else ""

    # Slides 2-N: Section slides
    for section in analysis.get("outline", []):
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = section.get("title", "")
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True

        # Summary
        summary_top = Inches(1.3)
        summary_box = slide.shapes.add_textbox(
            Inches(0.5), summary_top, Inches(9), Inches(1.5)
        )
        summary_frame = summary_box.text_frame
        summary_frame.text = section.get("summary", "")
        summary_frame.paragraphs[0].font.size = Pt(16)
        summary_frame.word_wrap = True

        # Screenshots (if any)
        screenshots = section.get("screenshots", [])
        if screenshots:
            img_top = Inches(3.0)
            img_width = Inches(4)
            img_left = Inches(3)  # Center horizontally

            for i, filename in enumerate(screenshots[:1]):  # Only first screenshot
                img_path = assets_dir / filename
                if img_path.exists():
                    slide.shapes.add_picture(
                        str(img_path), img_left, img_top, width=img_width
                    )

    # Final slide: Keywords summary
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "关键词汇总"
    tf = content.text_frame
    for kw in keywords:
        p = tf.add_paragraph()
        p.text = kw
        p.level = 0
        p.font.size = Pt(18)

    prs.save(str(output_path))
